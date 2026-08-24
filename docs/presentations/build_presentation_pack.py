#!/usr/bin/env python3
"""Build non-IT presentation pack for HIMA Automated Prooftest Reporting Solution."""

from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu
from reportlab.lib.colors import Color, HexColor, white, black
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.platypus import (
    SimpleDocTemplate,
    Paragraph,
    Spacer,
    Table,
    TableStyle,
    PageBreak,
    Image as RLImage,
    KeepTogether,
)

ROOT = Path(__file__).resolve().parent
DIAG = ROOT / "diagrams"
DIAG.mkdir(parents=True, exist_ok=True)

HIMA_BLUE = (0, 51, 102)
HIMA_BLUE_HEX = "#003366"
MID_BLUE = (0, 90, 158)
GREY = (74, 85, 104)
LIGHT = (232, 238, 244)
OPC_GREEN = (47, 111, 78)
SIL_AMBER = (180, 83, 9)
WHITE = (255, 255, 255)
DARK = (26, 32, 44)

VERSION = "1.77"
DATE = "2026-08-20"


def _font(size: int, bold: bool = False):
    candidates = [
        r"C:\Windows\Fonts\arialbd.ttf" if bold else r"C:\Windows\Fonts\arial.ttf",
        r"C:\Windows\Fonts\calibrib.ttf" if bold else r"C:\Windows\Fonts\calibri.ttf",
        r"C:\Windows\Fonts\segoeuib.ttf" if bold else r"C:\Windows\Fonts\segoeui.ttf",
    ]
    for path in candidates:
        if Path(path).is_file():
            return ImageFont.truetype(path, size)
    return ImageFont.load_default()


def _rounded(draw, box, fill, outline, width=3, radius=18):
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)


def _center_text(draw, box, text, font, fill=DARK):
    x0, y0, x1, y1 = box
    lines = text.split("\n")
    line_h = font.size + 4
    total = line_h * len(lines)
    y = y0 + (y1 - y0 - total) // 2
    for line in lines:
        bbox = draw.textbbox((0, 0), line, font=font)
        tw = bbox[2] - bbox[0]
        x = x0 + (x1 - x0 - tw) // 2
        draw.text((x, y), line, font=font, fill=fill)
        y += line_h


def _arrow(draw, start, end, color=GREY, width=4):
    draw.line([start, end], fill=color, width=width)
    x0, y0 = start
    x1, y1 = end
    # simple arrow head
    if abs(y1 - y0) > abs(x1 - x0):
        # vertical
        if y1 > y0:
            draw.polygon([(x1, y1), (x1 - 8, y1 - 14), (x1 + 8, y1 - 14)], fill=color)
        else:
            draw.polygon([(x1, y1), (x1 - 8, y1 + 14), (x1 + 8, y1 + 14)], fill=color)
    else:
        if x1 > x0:
            draw.polygon([(x1, y1), (x1 - 14, y1 - 8), (x1 - 14, y1 + 8)], fill=color)
        else:
            draw.polygon([(x1, y1), (x1 + 14, y1 - 8), (x1 + 14, y1 + 8)], fill=color)


def draw_architecture_png(path: Path) -> None:
    w, h = 1600, 1000
    img = Image.new("RGB", (w, h), WHITE)
    draw = ImageDraw.Draw(img)
    title_f = _font(36, True)
    box_f = _font(26, True)
    small_f = _font(20)
    legend_f = _font(18)

    draw.text((40, 24), "Architecture — who talks to whom", font=title_f, fill=HIMA_BLUE)

    # boxes
    sil = (80, 140, 420, 280)
    opc = (80, 360, 420, 500)
    tool = (560, 240, 1040, 420)
    db = (1160, 140, 1520, 280)
    rep = (1160, 320, 1520, 460)
    ui = (1160, 500, 1520, 640)
    plant = (80, 700, 420, 840)

    _rounded(draw, sil, (255, 244, 229), SIL_AMBER)
    _center_text(draw, sil, "SILworX\nengineering software", box_f)

    _rounded(draw, opc, (230, 244, 234), OPC_GREEN)
    _center_text(draw, opc, "OPC / live data link\n(X-OPC on station PC)", box_f)

    _rounded(draw, tool, LIGHT, HIMA_BLUE, width=5)
    _center_text(draw, tool, "Reporting tool\n(this solution — hub)", box_f, HIMA_BLUE)

    _rounded(draw, db, (243, 244, 246), GREY)
    _center_text(draw, db, "Database\nsaved results", box_f)

    _rounded(draw, rep, (243, 244, 246), GREY)
    _center_text(draw, rep, "Reports\nHTML + PDF", box_f)

    _rounded(draw, ui, (243, 244, 246), GREY)
    _center_text(draw, ui, "Local screen\nhttp://127.0.0.1:8080", box_f)

    _rounded(draw, plant, (243, 244, 246), GREY)
    _center_text(draw, plant, "Field devices\nHART proof tests", box_f)

    _arrow(draw, (420, 210), (560, 300), SIL_AMBER)
    draw.text((430, 175), "names & types\n(when connected)", font=small_f, fill=SIL_AMBER)

    _arrow(draw, (420, 430), (560, 360), OPC_GREEN)
    draw.text((430, 445), "live values +\nRunning (always)", font=small_f, fill=OPC_GREEN)

    _arrow(draw, (1040, 290), (1160, 210), GREY)
    _arrow(draw, (1040, 340), (1160, 390), GREY)
    _arrow(draw, (1040, 390), (1160, 560), GREY)

    _arrow(draw, (250, 700), (250, 500), GREY)
    draw.text((260, 600), "tests happen\nin plant / SILworX", font=small_f, fill=GREY)

    # legend
    draw.rectangle((80, 880, 1520, 970), fill=LIGHT, outline=HIMA_BLUE, width=2)
    draw.text((100, 895), "Legend:", font=legend_f, fill=HIMA_BLUE)
    draw.rectangle((220, 900, 260, 940), fill=(255, 244, 229), outline=SIL_AMBER, width=2)
    draw.text((270, 905), "SILworX", font=legend_f, fill=DARK)
    draw.rectangle((420, 900, 460, 940), fill=(230, 244, 234), outline=OPC_GREEN, width=2)
    draw.text((470, 905), "OPC live data", font=legend_f, fill=DARK)
    draw.rectangle((680, 900, 720, 940), fill=LIGHT, outline=HIMA_BLUE, width=3)
    draw.text((730, 905), "This reporting tool", font=legend_f, fill=DARK)
    draw.rectangle((1000, 900, 1040, 940), fill=(243, 244, 246), outline=GREY, width=2)
    draw.text((1050, 905), "Storage / screen", font=legend_f, fill=DARK)

    draw.text(
        (80, h - 28),
        "Caption: Live values always come from OPC. SILworX only helps identify devices when the tool is connected.",
        font=small_f,
        fill=GREY,
    )
    img.save(path)


def draw_flow_png(path: Path) -> None:
    w, h = 1400, 1600
    img = Image.new("RGB", (w, h), WHITE)
    draw = ImageDraw.Draw(img)
    title_f = _font(34, True)
    step_f = _font(22)
    draw.text((40, 24), "End-to-end flow — from test to report", font=title_f, fill=HIMA_BLUE)

    steps = [
        "1. Tool is running on the station PC",
        "2. Optional: Connect — attach to SILworX (names & types)",
        "3. If not connected — still find devices from OPC (carefully)",
        "4. Watch each device’s Running signal",
        "5. Running ON → note “test in progress”",
        "6. Running OFF → take results snapshot from OPC",
        "7. Save results in the database",
        "8. Create HTML + PDF report",
        "9. Operator opens report on the local screen",
    ]
    y = 90
    for i, text in enumerate(steps):
        box = (80, y, 900, y + 90)
        _rounded(draw, box, LIGHT, HIMA_BLUE, width=3, radius=14)
        _center_text(draw, (100, y, 880, y + 90), text, step_f, HIMA_BLUE)
        if i < len(steps) - 1:
            _arrow(draw, (490, y + 90), (490, y + 120), GREY, width=3)
        y += 120

    # branch
    branch = (960, 200, 1340, 420)
    _rounded(draw, branch, (255, 244, 229), SIL_AMBER, width=3, radius=14)
    _center_text(
        draw,
        branch,
        "Branch: Disconnect\nUnlinks this tool only.\nEngine keeps watching\nvia OPC. SILworX\nproject stays open.",
        step_f,
        DARK,
    )
    draw.line([(900, 245), (960, 280)], fill=SIL_AMBER, width=3)
    img.save(path)


def draw_connect_png(path: Path) -> None:
    w, h = 1600, 900
    img = Image.new("RGB", (w, h), WHITE)
    draw = ImageDraw.Draw(img)
    title_f = _font(34, True)
    head_f = _font(28, True)
    body_f = _font(22)
    draw.text((40, 24), "Connect vs Disconnect — what it really means", font=title_f, fill=HIMA_BLUE)

    left = (60, 100, 760, 720)
    right = (840, 100, 1540, 720)
    _rounded(draw, left, (230, 244, 234), OPC_GREEN, width=4, radius=20)
    _rounded(draw, right, (255, 244, 229), SIL_AMBER, width=4, radius=20)

    draw.text((100, 130), "CONNECT", font=head_f, fill=OPC_GREEN)
    for i, line in enumerate(
        [
            "• Links THIS reporting tool to SILworX",
            "• Reads which devices exist and their types",
            "• Does not take over SILworX",
            "• Does not open/close the engineer’s project",
        ]
    ):
        draw.text((100, 220 + i * 70), line, font=body_f, fill=DARK)

    draw.text((880, 130), "DISCONNECT", font=head_f, fill=SIL_AMBER)
    for i, line in enumerate(
        [
            "• Unlinks THIS tool only",
            "• Does NOT close the engineer’s project",
            "• Does NOT quit SILworX",
            "• Proof-test watching continues via OPC",
        ]
    ):
        draw.text((880, 220 + i * 70), line, font=body_f, fill=DARK)

    draw.text(
        (60, 780),
        "Think of Connect/Disconnect as plugging or unplugging this tool — not turning SILworX on or off.",
        font=body_f,
        fill=GREY,
    )
    img.save(path)


def draw_station_png(path: Path) -> None:
    w, h = 1400, 800
    img = Image.new("RGB", (w, h), WHITE)
    draw = ImageDraw.Draw(img)
    title_f = _font(34, True)
    box_f = _font(22, True)
    draw.text((40, 24), "One station PC — everything local", font=title_f, fill=HIMA_BLUE)

    pc = (200, 120, 1200, 680)
    _rounded(draw, pc, LIGHT, HIMA_BLUE, width=5, radius=24)
    draw.text((520, 150), "Windows station PC", font=_font(28, True), fill=HIMA_BLUE)

    items = [
        ((260, 240, 560, 360), "SILworX\n(may be open)", SIL_AMBER, (255, 244, 229)),
        ((640, 240, 940, 360), "X-OPC\nlive data link", OPC_GREEN, (230, 244, 234)),
        ((260, 420, 560, 540), "This reporting\ntool", HIMA_BLUE, LIGHT),
        ((640, 420, 940, 540), "Browser UI\n127.0.0.1:8080", GREY, (243, 244, 246)),
        ((980, 300, 1140, 480), "Reports\n+ Database\nfolders", GREY, (243, 244, 246)),
    ]
    for box, text, outline, fill in items:
        _rounded(draw, box, fill, outline, width=3, radius=14)
        _center_text(draw, box, text, box_f)
    img.save(path)


def write_drawio(path: Path, title: str, cells: list[str]) -> None:
    """Minimal editable draw.io XML wrapper."""
    body = "\n".join(cells)
    xml = f"""<?xml version="1.0" encoding="UTF-8"?>
<mxfile host="app.diagrams.net" modified="2026-08-20T00:00:00.000Z" agent="HIMA-presentation-pack" version="22.0.0">
  <diagram id="d1" name="{title}">
    <mxGraphModel dx="1200" dy="800" grid="1" gridSize="10" guides="1" tooltips="1" connect="1" arrows="1" fold="1" page="1" pageScale="1" pageWidth="1600" pageHeight="1000" math="0" shadow="0">
      <root>
        <mxCell id="0"/>
        <mxCell id="1" parent="0"/>
{body}
      </root>
    </mxGraphModel>
  </diagram>
</mxfile>
"""
    path.write_text(xml, encoding="utf-8")


def build_drawio_files() -> None:
    arch_cells = [
        '<mxCell id="sil" value="SILworX&#xa;names &amp; types" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#FFF4E5;strokeColor=#B45309;fontStyle=1" vertex="1" parent="1"><mxGeometry x="80" y="120" width="220" height="100" as="geometry"/></mxCell>',
        '<mxCell id="opc" value="OPC live data&#xa;values + Running" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#E6F4EA;strokeColor=#2F6F4E;fontStyle=1" vertex="1" parent="1"><mxGeometry x="80" y="280" width="220" height="100" as="geometry"/></mxCell>',
        '<mxCell id="tool" value="Reporting tool&#xa;(hub)" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#E8EEF4;strokeColor=#003366;strokeWidth=3;fontStyle=1" vertex="1" parent="1"><mxGeometry x="420" y="180" width="260" height="140" as="geometry"/></mxCell>',
        '<mxCell id="db" value="Database" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#F3F4F6;strokeColor=#4A5568" vertex="1" parent="1"><mxGeometry x="800" y="80" width="180" height="80" as="geometry"/></mxCell>',
        '<mxCell id="rep" value="Reports HTML/PDF" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#F3F4F6;strokeColor=#4A5568" vertex="1" parent="1"><mxGeometry x="800" y="200" width="180" height="80" as="geometry"/></mxCell>',
        '<mxCell id="ui" value="Local screen" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#F3F4F6;strokeColor=#4A5568" vertex="1" parent="1"><mxGeometry x="800" y="320" width="180" height="80" as="geometry"/></mxCell>',
        '<mxCell id="e1" style="endArrow=block;html=1;strokeColor=#B45309" edge="1" parent="1" source="sil" target="tool"><mxGeometry relative="1" as="geometry"/></mxCell>',
        '<mxCell id="e2" style="endArrow=block;html=1;strokeColor=#2F6F4E" edge="1" parent="1" source="opc" target="tool"><mxGeometry relative="1" as="geometry"/></mxCell>',
        '<mxCell id="e3" style="endArrow=block;html=1;strokeColor=#4A5568" edge="1" parent="1" source="tool" target="db"><mxGeometry relative="1" as="geometry"/></mxCell>',
        '<mxCell id="e4" style="endArrow=block;html=1;strokeColor=#4A5568" edge="1" parent="1" source="tool" target="rep"><mxGeometry relative="1" as="geometry"/></mxCell>',
        '<mxCell id="e5" style="endArrow=block;html=1;strokeColor=#4A5568" edge="1" parent="1" source="tool" target="ui"><mxGeometry relative="1" as="geometry"/></mxCell>',
    ]
    write_drawio(DIAG / "architecture.drawio", "Architecture", arch_cells)

    flow_cells = [
        f'<mxCell id="s{i}" value="{t}" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#E8EEF4;strokeColor=#003366" vertex="1" parent="1"><mxGeometry x="120" y="{80 + i * 70}" width="520" height="50" as="geometry"/></mxCell>'
        for i, t in enumerate(
            [
                "1 Tool running",
                "2 Optional Connect to SILworX",
                "3 Or find devices from OPC",
                "4 Watch Running",
                "5 Running ON",
                "6 Running OFF → snapshot",
                "7 Save database",
                "8 HTML + PDF",
                "9 Open on local screen",
            ]
        )
    ]
    for i in range(8):
        flow_cells.append(
            f'<mxCell id="fe{i}" style="endArrow=block;html=1;strokeColor=#4A5568" edge="1" parent="1" source="s{i}" target="s{i+1}"><mxGeometry relative="1" as="geometry"/></mxCell>'
        )
    write_drawio(DIAG / "end-to-end-flow.drawio", "End-to-end flow", flow_cells)

    cd_cells = [
        '<mxCell id="c" value="CONNECT&#xa;Link this tool to SILworX&#xa;Read names/types&#xa;Do not take over SILworX" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#E6F4EA;strokeColor=#2F6F4E;align=left;spacingLeft=10" vertex="1" parent="1"><mxGeometry x="60" y="120" width="360" height="200" as="geometry"/></mxCell>',
        '<mxCell id="d" value="DISCONNECT&#xa;Unlink this tool only&#xa;Do NOT close project&#xa;Do NOT quit SILworX&#xa;Watching continues via OPC" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#FFF4E5;strokeColor=#B45309;align=left;spacingLeft=10" vertex="1" parent="1"><mxGeometry x="480" y="120" width="360" height="200" as="geometry"/></mxCell>',
    ]
    write_drawio(DIAG / "connect-disconnect.drawio", "Connect Disconnect", cd_cells)


def _rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_run(shape, text, size=28, bold=False, color=HIMA_BLUE_HEX, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = _rgb(color)
    p.font.name = "Calibri"
    p.alignment = align


def add_bullets(shape, lines, size=20, color="#1A202C"):
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = _rgb(color)
        p.font.name = "Calibri"
        p.space_after = Pt(8)


def add_notes(slide, text: str) -> None:
    notes = slide.notes_slide
    notes.notes_text_frame.text = text.strip()


def build_pptx(pngs: dict[str, Path]) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def title_bar(slide, title: str, subtitle: str = ""):
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))  # rect
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0, 51, 102)
        box.line.fill.background()
        t = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.6))
        set_run(t, title, 28, True, "#FFFFFF")
        if subtitle:
            s = slide.shapes.add_textbox(Inches(0.4), Inches(0.75), Inches(12.5), Inches(0.35))
            set_run(s, subtitle, 14, False, "#D6E4F0")

    # SLIDE 1
    s = prs.slides.add_slide(blank)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0, 51, 102)
    bar.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(1.2))
    set_run(t, "HIMA Automated Prooftest Reporting Solution", 36, True, "#FFFFFF", PP_ALIGN.CENTER)
    t2 = s.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(11.5), Inches(0.8))
    set_run(t2, "Automatic proof-test records and reports from your station", 22, False, "#D6E4F0", PP_ALIGN.CENTER)
    t3 = s.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.8))
    set_run(
        t3,
        f"Version {VERSION}  ·  {DATE}  ·  For operations & engineering (non-IT overview)",
        16,
        False,
        "#A8C0D8",
        PP_ALIGN.CENTER,
    )
    add_notes(
        s,
        "Welcome. This is a plain-language overview of the HIMA Automated Prooftest Reporting Solution. "
        "It is for plant managers, engineers, operators, and auditors — not a software training. "
        "In about ten minutes we will cover what problem it solves, what it does and does not do, and how a normal day looks.",
    )

    # SLIDE 2
    s = prs.slides.add_slide(blank)
    title_bar(s, "The problem — before this tool")
    body = s.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5))
    add_bullets(
        body,
        [
            "Proof tests happen in the plant and in SILworX (the engineering software).",
            "Results are easy to lose, hard to collect in one place, and hard to show auditors.",
            "Manual copy/paste into folders or spreadsheets is slow and easy to get wrong.",
            "People spend time hunting files instead of reviewing the proof-test outcome.",
        ],
        24,
    )
    add_notes(
        s,
        "Start with pain. Engineers and operators already run proof tests. The gap is the paperwork trail: "
        "finding the right results, proving what happened, and preparing for audits without hunting through folders.",
    )

    # SLIDE 3
    s = prs.slides.add_slide(blank)
    title_bar(s, "The solution — in one sentence")
    quote = s.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.5), Inches(1.5))
    set_run(
        quote,
        "When a proof test finishes, the tool automatically saves the results and creates the report.",
        28,
        True,
        HIMA_BLUE_HEX,
        PP_ALIGN.CENTER,
    )
    note = s.shapes.add_textbox(Inches(1.2), Inches(4.5), Inches(10.8), Inches(1.5))
    set_run(
        note,
        "It must run on the same Windows PC as the OPC / X-OPC live data link "
        "(the live connection that carries proof-test values from the controller/PC).",
        18,
        False,
        "#4A5568",
        PP_ALIGN.CENTER,
    )
    add_notes(
        s,
        "Repeat the one sentence. Emphasize: the tool does not invent the test — it records the finish. "
        "Also say it must live on the station PC next to the live data link, not somewhere else on the network.",
    )

    # SLIDE 4
    s = prs.slides.add_slide(blank)
    title_bar(s, "What this tool does — and does not do")
    left = s.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(6), Inches(5.5))
    add_bullets(
        left,
        [
            "DOES",
            "• Watches devices on the station",
            "• Detects when a test starts and ends",
            "• Saves results to the database (a structured history store)",
            "• Creates HTML and PDF reports",
            "• Shows status on a local web screen",
        ],
        20,
    )
    right = s.shapes.add_textbox(Inches(7), Inches(1.4), Inches(5.8), Inches(5.5))
    add_bullets(
        right,
        [
            "DOES NOT",
            "• Run the HART proof test itself",
            "• Replace SILworX",
            "• Close the engineer’s SILworX project when you click Disconnect",
            "• Quit SILworX",
        ],
        20,
    )
    add_notes(
        s,
        "This is the most important clarity slide. SILworX and the devices run the proof test. "
        "This tool watches, stores, and reports. Disconnect never means ‘shut down SILworX’.",
    )

    # SLIDE 5
    s = prs.slides.add_slide(blank)
    title_bar(s, "Who uses it")
    rows = [
        ["Role", "What they use it for"],
        ["Operator / engineer", "See devices, connect the tool to SILworX, open reports"],
        ["Maintenance / safety", "Proof-test history and PDFs for compliance"],
        ["IT (optional)", "Install once on the station PC next to OPC"],
    ]
    table = s.shapes.add_table(len(rows), 2, Inches(0.8), Inches(1.8), Inches(11.5), Inches(3.5)).table
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(8)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(18)
                p.font.name = "Calibri"
                p.font.bold = r == 0
                p.font.color.rgb = RGBColor(255, 255, 255) if r == 0 else RGBColor(26, 32, 44)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
    add_notes(
        s,
        "Keep roles simple. Operators and engineers use the screen daily. Safety and maintenance care about the PDF trail. "
        "IT installs it once on the station — they are not the main audience for this talk.",
    )

    # SLIDE 6 architecture
    s = prs.slides.add_slide(blank)
    title_bar(s, "Big picture — architecture", "Live values always from OPC; SILworX helps with names/types when connected")
    s.shapes.add_picture(str(pngs["arch"]), Inches(0.4), Inches(1.25), width=Inches(12.5))
    add_notes(
        s,
        "Walk the diagram left to right. Field tests and SILworX sit on one side. OPC brings live values. "
        "The reporting tool is the hub. Database, reports, and the local screen sit on the right. "
        "Say the caption out loud once.",
    )

    # SLIDE 7 station
    s = prs.slides.add_slide(blank)
    title_bar(s, "Station layout — one Windows PC")
    s.shapes.add_picture(str(pngs["station"]), Inches(0.8), Inches(1.3), width=Inches(11.5))
    add_notes(
        s,
        "No network spaghetti. One station PC hosts SILworX if needed, X-OPC, this tool, the browser UI, "
        "and the reports/database folders. That is the whole footprint.",
    )

    # SLIDE 8 flow
    s = prs.slides.add_slide(blank)
    title_bar(s, "End-to-end flow — from test to report")
    s.shapes.add_picture(str(pngs["flow"]), Inches(1.8), Inches(1.15), height=Inches(6.1))
    add_notes(
        s,
        "Tell it as a story: tool running, optional Connect, watch Running, test ends, save, report, open. "
        "Point to the Disconnect branch: SILworX can stay open; watching continues on OPC.",
    )

    # SLIDE 9 connect
    s = prs.slides.add_slide(blank)
    title_bar(s, "Connect vs Disconnect — critical")
    s.shapes.add_picture(str(pngs["connect"]), Inches(0.4), Inches(1.25), width=Inches(12.5))
    add_notes(
        s,
        "Spend extra time here. Many people fear Disconnect. Reassure them: it only unplugs this tool. "
        "It does not close the engineer’s project and does not quit SILworX. Watching continues.",
    )

    # SLIDE 10 UI
    s = prs.slides.add_slide(blank)
    title_bar(s, "What you see on the screen")
    body = s.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
    add_bullets(
        body,
        [
            "Local page on that PC only (typically http://127.0.0.1:8080).",
            "You see devices, project, OPC server, alarms, and reports.",
            "Device table columns: Device, Type, OPC, Project, OPC server.",
            "The SILworX badge means “this tool is attached” — not “SILworX.exe is running”.",
            "Desktop shortcut (if installed): HIMA Prooftest Report.",
            "(Screenshot optional — wireframe: header buttons + device list + reports panel.)",
        ],
        22,
    )
    add_notes(
        s,
        "If you have a live station, show the browser briefly. Otherwise describe the layout. "
        "Clarify the badge wording so nobody confuses tool attachment with SILworX itself.",
    )

    # SLIDE 11 reports
    s = prs.slides.add_slide(blank)
    title_bar(s, "Reports and records")
    body = s.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
    add_bullets(
        body,
        [
            "Reports land under the station folder, for example:",
            "    C:\\HIMA Prooftest Reporting Tool\\… (configured reports path)",
            "HTML — quick view in a browser.",
            "PDF — archive and audit handout.",
            "The database keeps a history of saved results so you can look back later.",
        ],
        22,
    )
    add_notes(
        s,
        "Point to the station folder path. HTML for day-to-day viewing, PDF for formal archive. "
        "The database is simply the long-term memory of what was saved.",
    )

    # SLIDE 12 trust
    s = prs.slides.add_slide(blank)
    title_bar(s, "Safety and trust — short messages")
    body = s.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5))
    add_bullets(
        body,
        [
            "The tool does not change the proof-test logic in the PLC or SILworX.",
            "It records what the live data link (OPC) published when the test finished.",
            "It only accepts real proof-test result structures — not random signals.",
            "The screen is local-only by default — not for the open internet.",
        ],
        24,
    )
    add_notes(
        s,
        "This reassures safety and auditors. Recording is passive. Local-only reduces exposure. "
        "One plain phrase on structure filtering is enough — do not dive into software details.",
    )

    # SLIDE 13 day
    s = prs.slides.add_slide(blank)
    title_bar(s, "A typical day")
    body = s.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
    add_bullets(
        body,
        [
            "1. Engineer opens the SILworX project on the station.",
            "2. In the reporting tool, click Connect (this tool attaches).",
            "3. Operators run proof tests as usual.",
            "4. Reports appear automatically when each test finishes.",
            "5. Optional: Disconnect when engineering is done — SILworX can stay open; watching continues.",
        ],
        24,
    )
    add_notes(
        s,
        "Tell it as a five-beat story. Pause after Connect and after ‘reports appear automatically’. "
        "End with optional Disconnect so the fear topic is closed positively.",
    )

    # SLIDE 14 benefits
    s = prs.slides.add_slide(blank)
    title_bar(s, "Benefits")
    body = s.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
    add_bullets(
        body,
        [
            "Less manual paperwork after each proof test.",
            "The same record every time — consistent files and history.",
            "Faster answers when auditors ask “show me the last tests”.",
            "Clearer device list: the same tag in two projects shows as two lines "
            "(so you do not mix projects by accident).",
        ],
        24,
    )
    add_notes(
        s,
        "Keep benefits business-facing. The ‘same tag in two projects’ point is the only identity subtlety — "
        "explain it as two lines on the list, not as technical IDs.",
    )

    # SLIDE 15 FAQ
    s = prs.slides.add_slide(blank)
    title_bar(s, "FAQ — plain answers")
    body = s.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5))
    add_bullets(
        body,
        [
            "Does it run the test? No. SILworX and the devices do.",
            "Do we need SILworX connected always? No. OPC can continue; SILworX connection improves naming and types.",
            "Will Disconnect stop my tests? No.",
            "Where do I open it? Local station browser, or the desktop shortcut if installed.",
            "Is it for the whole plant network? By default it is only on that station PC.",
        ],
        22,
    )
    add_notes(
        s,
        "Use this if questions start early. Otherwise run it quickly. The Disconnect FAQ is the one people need most.",
    )

    # SLIDE 16 close
    s = prs.slides.add_slide(blank)
    title_bar(s, "Thank you — next steps")
    body = s.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.5))
    add_bullets(
        body,
        [
            f"Solution: HIMA Automated Prooftest Reporting  ·  version {VERSION}",
            "Questions: [project owner / station owner — fill in]",
            "Operators: see OPERATOR-CONNECT-DISCONNECT.md for Connect/Disconnect rules",
            "Handout: “HIMA Prooftest Reporting — at a glance” (one page)",
        ],
        22,
    )
    add_notes(
        s,
        "Leave the owner contact blank for the local team to fill. Offer the one-page handout and the operator note. "
        "Invite questions. Ten to twelve minutes plus discussion.",
    )

    out = ROOT / "HIMA-Prooftest-Reporting-Solution-Overview.pptx"
    prs.save(out)
    return out


def build_pdf_deck(pngs: dict[str, Path], pptx_path: Path) -> Path:
    """Printable landscape PDF mirroring the deck (text + key diagrams)."""
    out = ROOT / "HIMA-Prooftest-Reporting-Solution-Overview.pdf"
    doc = SimpleDocTemplate(
        str(out),
        pagesize=landscape(A4),
        leftMargin=18 * mm,
        rightMargin=18 * mm,
        topMargin=14 * mm,
        bottomMargin=14 * mm,
    )
    styles = getSampleStyleSheet()
    title = ParagraphStyle(
        "T",
        parent=styles["Heading1"],
        fontSize=22,
        textColor=HexColor(HIMA_BLUE_HEX),
        spaceAfter=10,
    )
    h = ParagraphStyle(
        "H",
        parent=styles["Heading2"],
        fontSize=16,
        textColor=HexColor(HIMA_BLUE_HEX),
        spaceBefore=6,
        spaceAfter=8,
    )
    body = ParagraphStyle("B", parent=styles["Normal"], fontSize=11, leading=15, textColor=HexColor("#1A202C"))
    story = []

    slides_text = [
        (
            "1. Title",
            f"<b>HIMA Automated Prooftest Reporting Solution</b><br/>"
            f"Automatic proof-test records and reports from your station<br/>"
            f"Version {VERSION} · {DATE} · For operations &amp; engineering (non-IT overview)",
        ),
        (
            "2. The problem",
            "Proof tests happen in plant / SILworX. Results are easy to lose and hard to show auditors. "
            "Manual copy/paste is slow and error-prone.",
        ),
        (
            "3. The solution",
            "<b>When a proof test finishes, the tool automatically saves the results and creates the report.</b><br/>"
            "Must run on the same Windows PC as the OPC / X-OPC live data link.",
        ),
        (
            "4. Does / Does not",
            "<b>Does:</b> watch devices, detect start/end, save database, create HTML/PDF, local screen.<br/>"
            "<b>Does not:</b> run the HART test, replace SILworX, close the engineer project on Disconnect, quit SILworX.",
        ),
        (
            "5. Who uses it",
            "Operator/engineer — devices, Connect, reports. Maintenance/safety — history and PDFs. "
            "IT (optional) — install on station PC next to OPC.",
        ),
        ("6. Architecture", "See diagram on next page / embedded below."),
        ("7. Station layout", "One Windows station PC: SILworX, X-OPC, this tool, browser UI, reports + database."),
        ("8. End-to-end flow", "See flow diagram. Optional Disconnect branch keeps OPC watching; SILworX stays open."),
        (
            "9. Connect vs Disconnect",
            "Connect links this tool only. Disconnect unlinks this tool only — never closes the project or quits SILworX.",
        ),
        (
            "10. Screen",
            "Local http://127.0.0.1:8080. Columns: Device, Type, OPC, Project, OPC server. "
            "Badge = tool attached, not “SILworX is running”.",
        ),
        (
            "11. Reports",
            r"Folder under C:\HIMA Prooftest Reporting Tool\… HTML for viewing, PDF for archive. Database keeps history.",
        ),
        (
            "12. Trust",
            "Does not change PLC/SILworX test logic. Records what OPC published. Local-only screen by default.",
        ),
        (
            "13. Typical day",
            "Open SILworX → Connect in tool → run tests → reports appear → optional Disconnect.",
        ),
        (
            "14. Benefits",
            "Less paperwork, consistent records, faster audits, clear list when the same tag exists in two projects.",
        ),
        (
            "15. FAQ",
            "Does not run the test. SILworX connect not always required. Disconnect does not stop tests. "
            "Open on local browser / desktop shortcut.",
        ),
        (
            "16. Close",
            f"Version {VERSION}. Questions: [project owner]. See OPERATOR-CONNECT-DISCONNECT.md.",
        ),
    ]

    for heading, text in slides_text:
        story.append(Paragraph(heading, h))
        story.append(Paragraph(text, body))
        if "Architecture" in heading:
            story.append(Spacer(1, 6))
            story.append(RLImage(str(pngs["arch"]), width=240 * mm, height=150 * mm))
        if "End-to-end" in heading:
            story.append(Spacer(1, 4))
            story.append(RLImage(str(pngs["flow"]), width=120 * mm, height=140 * mm))
        if "Connect vs" in heading:
            story.append(Spacer(1, 4))
            story.append(RLImage(str(pngs["connect"]), width=240 * mm, height=120 * mm))
        story.append(PageBreak())

    story.append(Paragraph("Source PowerPoint", h))
    story.append(Paragraph(str(pptx_path.name), body))
    doc.build(story)
    return out


def build_handout(pngs: dict[str, Path]) -> Path:
    out = ROOT / "HIMA-Prooftest-Reporting-at-a-glance.pdf"
    doc = SimpleDocTemplate(
        str(out),
        pagesize=A4,
        leftMargin=14 * mm,
        rightMargin=14 * mm,
        topMargin=12 * mm,
        bottomMargin=12 * mm,
    )
    styles = getSampleStyleSheet()
    title = ParagraphStyle(
        "HT",
        parent=styles["Heading1"],
        fontSize=16,
        textColor=HexColor(HIMA_BLUE_HEX),
        spaceAfter=6,
    )
    h = ParagraphStyle(
        "HH",
        parent=styles["Heading2"],
        fontSize=11,
        textColor=HexColor(HIMA_BLUE_HEX),
        spaceBefore=6,
        spaceAfter=3,
    )
    body = ParagraphStyle("HB", parent=styles["Normal"], fontSize=9, leading=12)
    story = [
        Paragraph("HIMA Prooftest Reporting — at a glance", title),
        Paragraph(
            f"Automatic proof-test records and reports from your station · Version {VERSION} · {DATE}",
            body,
        ),
        Paragraph("One sentence", h),
        Paragraph(
            "<b>When a proof test finishes, the tool automatically saves the results and creates the report.</b> "
            "SILworX/devices run the test; this tool watches, stores, and reports.",
            body,
        ),
        Paragraph("Does / Does not", h),
        Paragraph(
            "<b>Does:</b> watch devices, detect start/end, save results, create HTML/PDF, local screen.<br/>"
            "<b>Does not:</b> run the HART test, replace SILworX, close the engineer’s project on Disconnect, quit SILworX.",
            body,
        ),
        Paragraph("Connect vs Disconnect", h),
        Paragraph(
            "<b>Connect</b> = link this tool to SILworX for names/types. "
            "<b>Disconnect</b> = unlink this tool only. Watching continues via OPC (live data link).",
            body,
        ),
        Paragraph("Where to open", h),
        Paragraph(
            r"Local browser http://127.0.0.1:8080 on the station PC. Reports under C:\HIMA Prooftest Reporting Tool\…",
            body,
        ),
        Spacer(1, 4),
        RLImage(str(pngs["arch"]), width=180 * mm, height=112 * mm),
        Spacer(1, 4),
        Paragraph("Operators: see OPERATOR-CONNECT-DISCONNECT.md", body),
    ]
    doc.build(story)
    return out


def write_readme(paths: dict[str, Path]) -> None:
    text = f"""# Presentations — HIMA Automated Prooftest Reporting

| Field | Value |
|-------|--------|
| **Audience** | Non-IT: plant managers, process/safety engineers, operators, auditors |
| **Product truth** | Unified mode, code **{VERSION}** |
| **Talk length** | About 10–12 minutes |

## Files

| File | Purpose |
|------|---------|
| `{paths['pptx'].name}` | Main PowerPoint deck (speaker notes on every slide) |
| `{paths['pdf'].name}` | Printable PDF of the deck |
| `{paths['handout'].name}` | One-page handout |
| `diagrams/*.mmd` | Mermaid sources (editable) |
| `diagrams/*.drawio` | draw.io sources (editable) |
| `diagrams/*.png` | Diagram images embedded in the deck |

## How to open

1. Double-click the `.pptx` in PowerPoint or compatible software.
2. Use **Presenter View** to see speaker notes.
3. PDF deck / handout open in any PDF reader.
4. Edit diagrams in [diagrams.net](https://app.diagrams.net/) (open `.drawio`) or any Mermaid preview for `.mmd`.

## Message to protect

- This tool does **not** run the HART proof test.
- **Disconnect** only unlinks this tool — it does not close SILworX or the engineer’s project.
- Live values always come from OPC; SILworX helps with names/types when connected.

## Related operator note

`../OPERATOR-CONNECT-DISCONNECT.md`
"""
    (ROOT / "README.md").write_text(text, encoding="utf-8")


def main() -> None:
    print("Drawing PNGs…")
    pngs = {
        "arch": DIAG / "architecture.png",
        "flow": DIAG / "end-to-end-flow.png",
        "connect": DIAG / "connect-disconnect.png",
        "station": DIAG / "station-layout.png",
    }
    draw_architecture_png(pngs["arch"])
    draw_flow_png(pngs["flow"])
    draw_connect_png(pngs["connect"])
    draw_station_png(pngs["station"])
    print("Writing draw.io…")
    build_drawio_files()
    print("Building PPTX…")
    pptx_path = build_pptx(pngs)
    print("Building PDF deck…")
    pdf_path = build_pdf_deck(pngs, pptx_path)
    print("Building handout…")
    handout = build_handout(pngs)
    paths = {"pptx": pptx_path, "pdf": pdf_path, "handout": handout}
    write_readme(paths)
    print("DONE")
    for k, p in paths.items():
        print(f"  {k}: {p}")


if __name__ == "__main__":
    main()
