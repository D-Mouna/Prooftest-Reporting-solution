#!/usr/bin/env python3
"""Generate HIMA Prooftest architecture PowerPoint deck."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUTPUT = Path(__file__).resolve().parent / "HIMA-Prooftest-Architecture-Presentation.pptx"

SLIDES = [
    (
        "HIMA Automated Prooftest",
        [
            "How the tool works — architecture for developers & managers",
            "",
            "Audience: medium-level developers, project managers, integrators",
            "Spec: SPEC-001 v1.64 · Runtime UI: v1.80.54",
            "Web UI: http://127.0.0.1:8080/",
        ],
    ),
    (
        "What this tool does",
        [
            "Watches HART proof tests in SILworX via OPC (read-only)",
            "When a test ends, freezes all result values into SQL Server",
            "Generates HTML proof-test reports from frozen data",
            "Local web dashboard: devices, status, alarms, service control",
            "",
            "The tool attaches to SILworX — it does not open or close the engineer's project.",
        ],
    ),
    (
        "System context — external connections",
        [
            "SILworX REST API (HTTPS, ports 51710–51719) — device catalog, attach-only",
            "Session plugin WebSocket (8400–8409) — session ID for API attach",
            "HIMA X-OPC Classic DA — live .Running bit and result members (read-only)",
            "SQL Server (or SQLite fallback) — ProofTest_* tables, alarms, catalog",
            "Browser UI — REST polling on 127.0.0.1:8080",
            "",
            "Diagram: see Marp deck — 'System context' slide (Mermaid flowchart)",
        ],
    ),
    (
        "Layer architecture",
        [
            "Presentation: Graphic Interface (HTML/JS) + FastAPI controllers",
            "Application: ApplicationFacade — Catalog, LiveTest, Query, SILworX services",
            "Domain: pure logic — merge SILworX+OPC, Running edge detection",
            "Infrastructure: ProoftestService — poll loop, sync loop, health",
            "Annex adapters: OPC, API connexion, Plugin, Database, PDF generation",
            "",
            "Rule: UI never calls OPC or SQL directly — only through ApplicationFacade.",
            "Diagram: see Marp deck — 'Layer architecture' slide",
        ],
    ),
    (
        "Folder map",
        [
            "main.py + solution.ini — entry and configuration",
            "Graphic Interface/static/ — browser UI",
            "Annex codes/layers/presentation/ — REST API",
            "Annex codes/layers/application/facade.py — use cases",
            "Annex codes/layers/domain/ — merge and edge detection",
            "Tool Steps/service.py — engine host (poll + sync threads)",
            "Annex codes/OPC, API connexion, Plugin, Database, PDF — integrations",
            "Station data: C:\\HIMA Prooftest Reporting Tool\\",
        ],
    ),
    (
        "MVC-style UI request flow",
        [
            "User action in browser (View: app.js + HTML)",
            "POST/GET to FastAPI Controller",
            "Controller calls ApplicationFacade (never annex modules directly)",
            "Facade delegates to ProoftestService or domain services",
            "Service talks to OPC / SILworX / SQL via annex adapters",
            "JSON response → UI updates tiles, tables, buttons",
            "",
            "Diagram: see Marp deck — 'MVC-style request flow' sequence diagram",
        ],
    ),
    (
        "Connections — OPC (read-only)",
        [
            "Discover HIMA X-OPC servers (server_filter = HIMA.*)",
            "Browse OPC tree; bind {TAG}.Running and result member paths",
            "Poll every 1 second while engine is Running",
            "On test end: read all CSV-defined members once (snapshot)",
            "Never write or force OPC tags — safety constraint",
            "Modules: Annex codes/OPC/annex_opc.py, connection_opc.py",
        ],
    ),
    (
        "Connections — SILworX API + Plugin",
        [
            "Plugin WebSocket registers prooftest_session_plugin on port 8400+",
            "Receives TRIGGER_SESSION_ID_CHANGED when engineer opens/changes project",
            "REST API attaches to user-open project on port 51710+ (never open/local)",
            "Header: HIMA_SAPI_user_session_id from plugin cache",
            "Structure tree → Global Variables → device catalog merge with OPC",
            "Diagram: see Marp deck — 'SILworX API + Plugin' slide",
        ],
    ),
    (
        "Connections — Database & Web",
        [
            "SQL Server via ODBC; SQLite fallback if SQL unavailable",
            "ProofTest_* tables created from Results Structure CSV (sync_schema_case2)",
            "Schema sync runs once per process (avoids Stop→Start hang)",
            "FastAPI serves UI at 127.0.0.1:8080; optional token auth",
            "Destructive API calls restricted to localhost client",
        ],
    ),
    (
        "Configuration — solution.ini",
        [
            "[Paths] — station root, report folders, Results Structures",
            "[Database] — SQL instance and catalog database name",
            "[SILworX] — API ports, plugin name, host",
            "[OPC] — poll interval, shape gate, server filter",
            "[Web] — port, auth token, localhost bypass",
            "[Service] — auto-start, sync_triggers list",
            "Single config file — restart service after changes",
        ],
    ),
    (
        "Rules & safety constraints",
        [
            "OPC read-only — cannot disturb running safety logic",
            "Attach-only SILworX — engineer controls project open/close",
            "Disconnect detaches tool only — does not quit SILworX or kill c3.exe",
            "Release SILworX — special uninstall path; then OPC-only until Re-integrate",
            "Localhost-only: Start, Stop, Connect, Archive, alarm reset",
            "Freeze at test end: SQL insert before report queue",
            "Operator detach: no auto-reconnect after Disconnect until Connect",
            "Windows only — OPC Classic DA requirement",
        ],
    ),
    (
        "Decisions — engine lifecycle",
        [
            "Stopped → Starting → Running (POST /api/start)",
            "Running: poll-loop 1s, sync-loop 0.5s, report worker thread",
            "Running → Stopping → Stopped (POST /api/stop; UI stays open)",
            "Shutdown exits entire process (POST /api/shutdown or stop_service.ps1)",
            "",
            "Diagram: see Marp deck — 'engine lifecycle' state diagram",
        ],
    ),
    (
        "Decisions — proof test detection",
        [
            "Every 1s: read .Running with Good OPC quality",
            "RunningEdgeDetector: 0→1 = test started; 1→0 = test ended",
            "On end: collect full OPC snapshot → INSERT ProofTest_* row",
            "Enqueue HTML report job; worker uses frozen SQL row (no re-read)",
            "Bad quality / flicker → alarm S5, test marked interrupted",
            "",
            "Diagram: see Marp deck — 'proof test detection' flowchart",
        ],
    ),
    (
        "Decisions — background sync",
        [
            "Every 2s: check SILworX install, operator detach, open projects",
            "G-11: auto-release API if SILworX uninstalled → OPC-only",
            "Sync triggers: silworx_session, code_generation, download, results_structures",
            "Triggers refresh device catalog (SILworX merge + OPC bind)",
            "Disconnect sets _operator_detached — sync will not auto-resume API",
        ],
    ),
    (
        "Decisions — health & alarms",
        [
            "GET /api/health — engine, OPC, SILworX, plugins, device counts",
            "2s cache when Running; fast path when Stopped",
            "Alarms S1–S7 map to setup and runtime phases",
            "G-11 alarm when SILworX released for uninstall",
            "Alarms persisted in DB; UI polls /api/alarms every 2s",
        ],
    ),
    (
        "User interface — four pages",
        [
            "Monitor — device list, reports, archive tools",
            "Status — health tiles (OPC, SILworX, DB, plugin sessions)",
            "Alarms — error list, acknowledge, reset",
            "Service — Start/Stop, Connect/Disconnect, Release/Re-integrate SILworX",
            "Sidebar: running tests, history, theme",
            "Top bar chips: Devices, OPC count, Service, SILworX attach",
        ],
    ),
    (
        "UI — polling architecture",
        [
            "No WebSocket to browser — simple REST poll every 2 seconds",
            "pollStatus(): /api/health → /api/alarms → /api/running-tests → /api/devices",
            "User mutations: POST to localhost-only endpoints",
            "Manual Refresh: POST /api/refresh (catalog rebuild, up to 180s)",
            "",
            "Diagram: see Marp deck — 'UI architecture' slide",
        ],
    ),
    (
        "Main API endpoints",
        [
            "GET /api/health, /api/devices, /api/reports, /api/alarms",
            "POST /api/start, /api/stop (localhost)",
            "POST /api/silworx/connect, disconnect, release, reintegrate",
            "POST /api/refresh, /api/archives/*, /api/devices/keep-opc",
            "GET /api/reports/open — render HTML report in browser",
            "Full map: Annex codes/layers/presentation/controllers.py",
        ],
    ),
    (
        "Deployment & operations",
        [
            "Normal: SILworX API + OPC when both available",
            "OPC fallback when SILworX down or after Disconnect",
            "Release SILworX → OPC-only until Re-integrate",
            "Start: run_service.ps1 or Windows auto-start task",
            "Stop engine: Service page (UI stays up)",
            "Full exit: stop_service.ps1 before SILworX uninstall",
        ],
    ),
    (
        "Typical engineer workflow",
        [
            "Engineer opens SILworX project and runs proof test",
            "Tool attaches API + polls OPC .Running",
            "Test end → SQL freeze + HTML report",
            "Engineer opens report from Monitor page in browser",
            "Engineer workflow unchanged — tool runs beside SILworX",
            "",
            "Diagram: see Marp deck — 'Typical day' sequence diagram",
        ],
    ),
    (
        "Summary — three takeaways",
        [
            "1. Layers: UI → Facade → Service → Annex (never skip layers)",
            "2. Connections: OPC read-only + SILworX attach + plugin WS + SQL freeze",
            "3. Decisions: .Running edge, sync triggers, localhost guards",
            "",
            "Mermaid diagrams: open HIMA-Prooftest-Architecture-Presentation.md",
            "Export via Marp for VS Code → PPTX with diagrams embedded",
        ],
    ),
    (
        "Questions?",
        [
            "Spec: Report Solution/Specifications/SPEC-001-v1.64-…",
            "Code: Codes/HIMA-Prooftest-Solution-Current/",
            "Tests: Annex codes/Tool test/test_step13_hardening.py",
        ],
    ),
]


def add_title_slide(prs: Presentation, title: str, lines: list[str]) -> None:
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(lines):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18 if line else 12)
        if line.startswith("Diagram:"):
            p.font.italic = True
            p.font.color.rgb = RGBColor(0x00, 0x55, 0x99)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = "HIMA Automated Prooftest"
    slide.placeholders[1].text = "Architecture overview for developers & managers\nSPEC-001 · 2026"

    for title, bullets in SLIDES[1:]:
        add_title_slide(prs, title, bullets)

    prs.save(str(OUTPUT))
    print(f"Wrote {OUTPUT}")


if __name__ == "__main__":
    main()
